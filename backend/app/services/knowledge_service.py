import asyncio
import base64
import logging
import io
from uuid import UUID
import pdfplumber
from docx import Document as DocxDocument
from pptx import Presentation
import pandas as pd
from groq import Groq
from app.db.supabase import get_supabase
from app.config import settings

logger = logging.getLogger(__name__)

_MAX_TEXT_CHARS = 50_000

# Chunking: ~1500 chars (~400-450 tokens) stays under Voyage's per-input limit and
# keeps each chunk semantically tight. 200-char overlap preserves cross-boundary context.
_CHUNK_CHARS = 1500
_CHUNK_OVERLAP = 200
_MATCH_COUNT = 5


def _chunk_text(text: str) -> list[str]:
    """Split text into overlapping windows, preferring paragraph/sentence boundaries."""
    text = text.strip()
    if not text:
        return []
    if len(text) <= _CHUNK_CHARS:
        return [text]
    chunks: list[str] = []
    start = 0
    n = len(text)
    while start < n:
        end = min(start + _CHUNK_CHARS, n)
        if end < n:
            window = text[start:end]
            for sep in ("\n\n", "\n", ". "):
                idx = window.rfind(sep)
                if idx > _CHUNK_CHARS // 2:
                    end = start + idx + len(sep)
                    break
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        if end >= n:
            break
        start = max(end - _CHUNK_OVERLAP, start + 1)
    return chunks


async def _index_chunks(document_id: UUID, tenant_id: str, text: str, db) -> int:
    """Chunk → embed → replace this document's chunks. Returns chunk count."""
    from app.services.embeddings import embed_texts, to_pgvector

    chunks = _chunk_text(text)
    if not chunks:
        return 0
    embeddings = await embed_texts(chunks, input_type="document")
    if len(embeddings) != len(chunks):
        raise ValueError(f"Embedding count {len(embeddings)} != chunk count {len(chunks)}")

    db.table("knowledge_chunks").delete().eq("document_id", str(document_id)).execute()
    for idx, (chunk, emb) in enumerate(zip(chunks, embeddings)):
        db.rpc(
            "insert_knowledge_chunk",
            {
                "p_tenant_id": tenant_id,
                "p_document_id": str(document_id),
                "p_chunk_index": idx,
                "p_content": chunk,
                "p_embedding": to_pgvector(emb),
            },
        ).execute()
    return len(chunks)

_groq_client = Groq(api_key=settings.groq_api_key) if settings.groq_api_key else None
_VISION_MODEL = "llama-3.2-11b-vision-preview"


def extract_text_from_file(file_content: bytes, filename: str, mime_type: str) -> str:
    file_obj = io.BytesIO(file_content)
    text = ""
    try:
        if mime_type == "application/pdf" or filename.endswith(".pdf"):
            with pdfplumber.open(file_obj) as pdf:
                text = "\n".join(page.extract_text() or "" for page in pdf.pages)

        elif mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document" or filename.endswith(".docx"):
            doc = DocxDocument(file_obj)
            text = "\n".join(p.text for p in doc.paragraphs)

        elif mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation" or filename.endswith(".pptx"):
            prs = Presentation(file_obj)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"

        elif mime_type == "text/csv" or filename.endswith(".csv"):
            df = pd.read_csv(file_obj)
            text = df.to_string()

        elif mime_type in ("application/vnd.ms-excel", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet") or filename.endswith((".xls", ".xlsx")):
            df = pd.read_excel(file_obj)
            text = df.to_string()

        elif mime_type.startswith("image/"):
            if not _groq_client:
                raise ValueError("GROQ_API_KEY not configured — cannot extract text from images")
            b64 = base64.b64encode(file_content).decode("utf-8")
            response = _groq_client.chat.completions.create(
                model=_VISION_MODEL,
                messages=[{
                    "role": "user",
                    "content": [
                        {"type": "text", "text": "Extract all text from this image. Return only the extracted text, no commentary."},
                        {"type": "image_url", "image_url": {"url": f"data:{mime_type};base64,{b64}"}},
                    ],
                }],
                temperature=0.0,
                max_tokens=4000,
            )
            text = (response.choices[0].message.content or "").strip()

        elif mime_type == "text/plain" or filename.endswith(".txt"):
            text = file_content.decode("utf-8")

        else:
            try:
                text = file_content.decode("utf-8")
            except Exception:
                text = ""

    except Exception as e:
        logger.error(f"Text extraction failed for {filename}: {e}")
        raise ValueError(f"Could not extract text from {filename}: {e}")

    return text.strip()


async def process_document(document_id: UUID, tenant_id: str, file_content: bytes, filename: str, mime_type: str):
    db = get_supabase()
    try:
        text = await asyncio.to_thread(extract_text_from_file, file_content, filename, mime_type)
        if not text:
            raise ValueError("No text could be extracted from the file.")

        # full_text is retained as the embedding-failure fallback for retrieval.
        db.table("knowledge_documents").update({
            "status": "indexed",
            "full_text": text[:_MAX_TEXT_CHARS],
        }).eq("id", str(document_id)).execute()

        try:
            chunk_count = await _index_chunks(document_id, tenant_id, text[:_MAX_TEXT_CHARS], db)
            logger.info(f"Document {document_id} indexed — {len(text)} chars, {chunk_count} chunks embedded")
        except Exception as embed_err:
            # Indexing succeeded for full-text fallback; embeddings failed. Don't fail the
            # whole upload — retrieval falls back to full_text injection until re-indexed.
            logger.error(f"Chunk embedding failed for {document_id}: {embed_err}. Full-text fallback active.")

    except Exception as e:
        logger.error(f"Document processing failed for {document_id}: {e}")
        db.table("knowledge_documents").update({
            "status": "failed",
            "error_message": str(e),
        }).eq("id", str(document_id)).execute()


def _full_text_context(tenant_id: str, db) -> str:
    """Legacy full-text injection — used as the embedding-failure fallback."""
    res = (
        db.table("knowledge_documents")
        .select("name,full_text")
        .eq("tenant_id", tenant_id)
        .eq("status", "indexed")
        .execute()
    )
    parts = [
        f"=== {doc['name']} ===\n{doc['full_text']}"
        for doc in (res.data or [])
        if doc.get("full_text")
    ]
    return "\n\n".join(parts)


def _format_excerpts(rows: list[dict]) -> str:
    return "\n\n".join(f"=== excerpt {i + 1} ===\n{r['content']}" for i, r in enumerate(rows))


async def _semantic_search(db, tenant_id: str, query: str) -> list[dict]:
    """Vector similarity via Jina embeddings + match_knowledge_chunks RPC."""
    from app.services.embeddings import embed_query, to_pgvector

    q_emb = await embed_query(query)
    res = db.rpc(
        "match_knowledge_chunks",
        {"query_embedding": to_pgvector(q_emb), "p_tenant_id": tenant_id, "match_count": _MATCH_COUNT},
    ).execute()
    return res.data or []


def _keyword_search(db, tenant_id: str, query: str) -> list[dict]:
    """Postgres full-text + trigram match — no external API, language-agnostic tokens."""
    res = db.rpc(
        "keyword_match_chunks",
        {"p_query": query, "p_tenant_id": tenant_id, "match_count": _MATCH_COUNT},
    ).execute()
    return res.data or []


def _rrf_merge(lists: list[list[dict]], top_n: int, k: int = 60) -> list[dict]:
    """Reciprocal Rank Fusion: blend ranked result lists by content, best-first."""
    scores: dict[str, float] = {}
    keep: dict[str, dict] = {}
    for lst in lists:
        for rank, row in enumerate(lst):
            key = row.get("content") or ""
            if not key:
                continue
            scores[key] = scores.get(key, 0.0) + 1.0 / (k + rank + 1)
            keep[key] = row
    ordered = sorted(scores, key=lambda c: scores[c], reverse=True)[:top_n]
    return [keep[c] for c in ordered]


async def get_knowledge_context(tenant_id: str, query: str | None = None) -> str:
    """
    Retrieve the most relevant knowledge-base excerpts for this tenant's message.

    Retrieval mode is per-tenant via the `kb_retrieval_mode` setting:
      - "semantic" (default): Jina vector search — meaning-based, multilingual.
      - "keyword": Postgres full-text/trigram — exact tokens, no embedding API.
      - "hybrid": both, fused with Reciprocal Rank Fusion.

    Always falls back to full-text injection when the query is empty, retrieval errors,
    or nothing matches — North Star: a provider hiccup must never blank the knowledge base.
    """
    db = get_supabase()
    if not query or not query.strip():
        try:
            return _full_text_context(tenant_id, db)
        except Exception as e:
            logger.error(f"Knowledge full-text fetch failed: {e}")
            return ""

    from app.config_dynamic import get_setting

    mode = (get_setting("kb_retrieval_mode", fallback="semantic", tenant_id=tenant_id) or "semantic").lower()

    rows: list[dict] = []
    try:
        if mode == "keyword":
            rows = _keyword_search(db, tenant_id, query)
        elif mode == "hybrid":
            sem: list[dict] = []
            kw: list[dict] = []
            try:
                sem = await _semantic_search(db, tenant_id, query)
            except Exception as e:
                logger.warning(f"Hybrid: semantic leg failed for tenant {tenant_id}: {e}")
            try:
                kw = _keyword_search(db, tenant_id, query)
            except Exception as e:
                logger.warning(f"Hybrid: keyword leg failed for tenant {tenant_id}: {e}")
            rows = _rrf_merge([sem, kw], _MATCH_COUNT)
        else:
            rows = await _semantic_search(db, tenant_id, query)
    except Exception as e:
        logger.warning(f"KB retrieval ({mode}) failed for tenant {tenant_id}: {e}. Falling back to full-text.")

    if rows:
        return _format_excerpts(rows)
    logger.info(f"No chunks matched ({mode}) for tenant {tenant_id} — falling back to full-text injection")

    try:
        return _full_text_context(tenant_id, db)
    except Exception as e:
        logger.error(f"Knowledge full-text fallback failed: {e}")
        return ""


async def reindex_tenant(tenant_id: str) -> dict:
    """Re-chunk + embed every indexed document for a tenant (backfill for pre-RAG docs)."""
    db = get_supabase()
    res = (
        db.table("knowledge_documents")
        .select("id,full_text")
        .eq("tenant_id", tenant_id)
        .eq("status", "indexed")
        .execute()
    )
    docs = res.data or []
    total_chunks = 0
    reindexed = 0
    for doc in docs:
        text = doc.get("full_text")
        if not text:
            continue
        try:
            total_chunks += await _index_chunks(UUID(doc["id"]), tenant_id, text, db)
            reindexed += 1
        except Exception as e:
            logger.error(f"Reindex failed for document {doc['id']}: {e}")
    return {"documents_reindexed": reindexed, "chunks_embedded": total_chunks}
